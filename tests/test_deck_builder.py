"""Tests for DeckBuilder — Wk6 polish coverage."""
from backend.tools.deck_builder import DeckBuilder


def test_narrative_to_bullets_splits_sentences():
    builder = DeckBuilder()
    result = builder._narrative_to_bullets("First point. Second point. Third point.")
    assert result == ["First point.", "Second point.", "Third point."]


def test_narrative_to_bullets_handles_empty_string():
    builder = DeckBuilder()
    result = builder._narrative_to_bullets("")
    assert result == ["Narrative not available."]


def test_narrative_to_bullets_handles_none():
    builder = DeckBuilder()
    result = builder._narrative_to_bullets(None)
    assert result == ["Narrative not available."]


def test_narrative_to_bullets_strips_newlines():
    builder = DeckBuilder()
    result = builder._narrative_to_bullets("Line one.\nLine two.")
    assert result == ["Line one.", "Line two."]


def test_build_deck_includes_competitors_slide(tmp_path):
    builder = DeckBuilder()
    ceo_output = {
        "idea": "Test idea",
        "ceo_narrative": "We solve X.",
        "cmo_output": {
            "market": {"tam": "$1B"},
            "competitors": [{"name": "CompA", "summary": "does X"}],
        },
        "cfo_output": {},
    }
    output_path = str(tmp_path / "test.pptx")
    builder.build_deck(ceo_output, output_path)

    from pptx import Presentation
    prs = Presentation(output_path)
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "Competitive Landscape" in titles


def test_build_deck_includes_product_tech_slide_when_cto_present(tmp_path):
    builder = DeckBuilder()
    ceo_output = {
        "idea": "Test idea",
        "ceo_narrative": "We solve X.",
        "cmo_output": {},
        "cfo_output": {},
        "cto_output": {
            "mvp_features": [{"name": "Feature A"}],
            "tech_stack": {"frontend": "React", "backend": "FastAPI", "database": "Postgres"},
        },
    }
    output_path = str(tmp_path / "test2.pptx")
    builder.build_deck(ceo_output, output_path)

    from pptx import Presentation
    prs = Presentation(output_path)
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    assert "Product & Tech" in titles


def test_build_deck_handles_no_competitors_gracefully(tmp_path):
    builder = DeckBuilder()
    ceo_output = {
        "idea": "Test idea",
        "ceo_narrative": "We solve X.",
        "cmo_output": {"competitors": []},
        "cfo_output": {},
    }
    output_path = str(tmp_path / "test3.pptx")
    path = builder.build_deck(ceo_output, output_path)
    assert path == output_path  # no crash


def test_build_deck_produces_six_slides(tmp_path):
    builder = DeckBuilder()
    ceo_output = {
        "idea": "Test idea",
        "ceo_narrative": "We solve X. We do Y.",
        "cmo_output": {"market": {"tam": "$1B"}, "competitors": [{"name": "A", "summary": "s"}]},
        "cfo_output": {"funding_ask": {"amount": "$100k"}},
        "cto_output": {"mvp_features": [{"name": "F1"}], "tech_stack": {}},
    }
    output_path = str(tmp_path / "test4.pptx")
    builder.build_deck(ceo_output, output_path)

    from pptx import Presentation
    prs = Presentation(output_path)
    assert len(prs.slides._sldIdLst) == 6