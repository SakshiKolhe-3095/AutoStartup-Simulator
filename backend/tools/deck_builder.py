"""
Deck Builder — auto-generates pitch deck slides using python-pptx.
Owner: Sakshi
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


TITLE_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(18)

ACCENT_COLOR = RGBColor(0x2E, 0x7D, 0x32)      # dark green — matches frontend accent
TITLE_TEXT_COLOR = RGBColor(0x1B, 0x1B, 0x1B)  # near-black
BODY_TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)


class DeckBuilder:
    def __init__(self, template_path: str = None):
        self.template_path = template_path

    @staticmethod
    def _narrative_to_bullets(narrative: str) -> list:
        """Split a narrative paragraph into clean sentence-level bullets."""
        if not narrative:
            return ["Narrative not available."]
        raw_sentences = narrative.replace("\n", " ").split(". ")
        bullets = [s.strip().rstrip(".") + "." for s in raw_sentences if s.strip()]
        return bullets or ["Narrative not available."]

    @staticmethod
    def _set_body_text(slide, lines: list):
        """Write a list of lines into the body placeholder with consistent styling."""
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(lines):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = line
            for run in p.runs:
                run.font.size = BODY_FONT_SIZE
                run.font.color.rgb = BODY_TEXT_COLOR

    def _add_slide(self, prs, title: str, lines: list):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
        title_run.font.size = TITLE_FONT_SIZE
        title_run.font.bold = True
        title_run.font.color.rgb = ACCENT_COLOR
        self._set_body_text(slide, lines)
        return slide

    def build_deck(self, ceo_output: dict, output_path: str = "output_deck.pptx"):
        """
        Auto-populates slides from CEO-agent output.
        ceo_output expected keys: idea, ceo_narrative, cmo_output, cfo_output, cto_output
        """
        prs = Presentation(self.template_path) if self.template_path else Presentation()
        idea = ceo_output.get("idea", "Untitled Startup")
        narrative = ceo_output.get("ceo_narrative", "")
        cmo = ceo_output.get("cmo_output", {}) or {}
        cfo = ceo_output.get("cfo_output", {}) or {}
        cto = ceo_output.get("cto_output", {}) or {}

        # Slide 1: Title — accent color, bold headline
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = idea
        title_run = title_shape.text_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(40)
        title_run.font.bold = True
        title_run.font.color.rgb = ACCENT_COLOR

        subtitle = slide.placeholders[1]
        subtitle.text = "AutoStartup Simulator — Pitch Deck"
        subtitle_run = subtitle.text_frame.paragraphs[0].runs[0]
        subtitle_run.font.color.rgb = BODY_TEXT_COLOR
        subtitle_run.font.italic = True

        # Slide 2: Pitch narrative
        self._add_slide(prs, "The Pitch", self._narrative_to_bullets(narrative))

        # Slide 3: Market (from CMO)
        market = cmo.get("market", {})
        self._add_slide(prs, "Market Opportunity", [
            f"TAM: {market.get('tam', 'N/A')}",
            f"SAM: {market.get('sam', 'N/A')}",
            f"SOM: {market.get('som', 'N/A')}",
        ])

        # Slide 4: Competitors (from CMO)
        competitors = cmo.get("competitors", [])
        if competitors:
            lines = [f"{c.get('name', 'Unknown')}: {c.get('summary', '')}" for c in competitors[:5]]
        else:
            lines = ["No competitor data available."]
        self._add_slide(prs, "Competitive Landscape", lines)

        # Slide 5: Product / Tech Stack (from CTO)
        mvp_features = cto.get("mvp_features", [])
        tech_stack = cto.get("tech_stack", {})
        product_lines = []
        if mvp_features:
            product_lines.append("MVP Features:")
            product_lines += [f"- {f.get('name', '')}" for f in mvp_features[:5]]
        if tech_stack:
            product_lines.append(
                f"Stack: {tech_stack.get('frontend', 'N/A')} / "
                f"{tech_stack.get('backend', 'N/A')} / "
                f"{tech_stack.get('database', 'N/A')}"
            )
        if not product_lines:
            product_lines = ["Product details not available."]
        self._add_slide(prs, "Product & Tech", product_lines)

        # Slide 6: Financials (from CFO)
        funding = cfo.get("funding_ask", {})
        self._add_slide(prs, "Financials", [
            f"Funding ask: {funding.get('amount', 'N/A')}",
            f"Use of funds: {funding.get('use_of_funds', 'N/A')}",
        ])

        prs.save(output_path)
        return output_path


if __name__ == "__main__":
    builder = DeckBuilder()
    sample_output = {
        "idea": "AI note-taking app",
        "ceo_narrative": "We help busy professionals capture ideas instantly.",
        "cmo_output": {
            "market": {"tam": "$1B", "sam": "$100M", "som": "$5M"},
            "competitors": [{"name": "Notion", "summary": "General workspace tool"}],
        },
        "cfo_output": {"funding_ask": {"amount": "$250k", "use_of_funds": "hiring + infra"}},
        "cto_output": {
            "mvp_features": [{"name": "Voice-to-text capture"}],
            "tech_stack": {"frontend": "React", "backend": "FastAPI", "database": "Postgres"},
        },
    }
    path = builder.build_deck(sample_output)
    print(f"Deck saved: {path}")